import pptx
import os

def list_all_textboxes_with_names(pptx_path):
    presentation = pptx.Presentation(pptx_path)
    
    for slide_num, slide in enumerate(presentation.slides, start=1):
        print(f"\n🔹 Slide {slide_num}:")
        count = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                count += 1
                text_preview = shape.text.strip().replace("\n", " ")[:50]
                print(f"  [{count}] Shape Name: '{shape.name}' | Text: '{text_preview}'")

# Example usage
pptx_path = os.path.join("powerpoints", "template2.pptx")
list_all_textboxes_with_names(pptx_path)
