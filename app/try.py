import pptx
import os
import requests
import json

presentation = pptx.Presentation(os.path.join("powerpoints", "template3.pptx"))

def list_text_boxes(presentation, slide_num):
    slide = presentation.slides[slide_num-1]
    text_boxes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            text_boxes.append(shape.text)
    return text_boxes

for idx, text in enumerate(list_text_boxes(presentation,9), 1):
   print(f"Text Box {idx}: {text}")
