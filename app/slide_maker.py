import os
import pptx
from datetime import datetime
from pptx.dml.color import RGBColor


def slide_generation_1(arr, template_type):
    # Load the base presentation
    presentation = pptx.Presentation(os.path.join("powerpoints", f"{template_type}.pptx"))

    def update_text_of_textbox(presentation, slide_num, text_box_id, new_text):
        slide = presentation.slides[slide_num - 1]
        count = 0
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                count += 1
                if count == text_box_id:
                    text_frame = shape.text_frame
                    print(f"[Slide {slide_num}] Box {count}: '{shape.text.strip()}'")
                    first_paragraph = text_frame.paragraphs[0]
                    first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()

                    # Preserve formatting
                    font = first_run.font
                    font_name, font_size = font.name, font.size
                    font_bold, font_italic = font.bold, font.italic
                    # Save existing font styles if needed
                    font_underline = font.underline
                    try:
                        font_color = font.color.rgb
                    except:
                        font_color = None

                    try:
                        # Clear and set new text
                        text_frame.clear()
                        new_run = text_frame.paragraphs[0].add_run()
                        if not isinstance(new_text, str):
                            new_text = str(new_text)
                        new_run.text = new_text
                    
                        print(f"📝 Updating Slide {slide_num}, Box {text_box_id} → Text: {new_text[:30]}...")

                        # Reapply formatting
                        new_run.font.name = font_name
                        new_run.font.size = font_size
                        new_run.font.bold = font_bold
                        new_run.font.italic = font_italic
                        new_run.font.underline = font_underline
                        if font_color:
                            try:
                                if isinstance(font_color, RGBColor):
                                    new_run.font.color.rgb = font_color
                                elif isinstance(font_color, tuple) and len(font_color) == 3:
                                    new_run.font.color.rgb = RGBColor(*font_color)
                            except Exception as e:
                                print(f"⚠️ Font color issue: {e}")
                    except Exception as e:
                        print(f"❌ Failed to update text: {e}")        
                    return

    # Define textbox mappings for each template
    template_mappings = {
        "template1": {1: [1, 3], **{i: [3, 4] for i in range(2, 10)}},
        "template3": {1: [2, 3], **{i: [1, 2] for i in range(2, 10)}}
    }

    if template_type not in template_mappings:
        raise ValueError(f"❌ Unknown template type: {template_type}")

    mapping = template_mappings[template_type]

    # Fill slides based on the selected template mapping
    for slide_num in range(1, min(len(arr), 9) + 1):
        text_box_ids = mapping.get(slide_num, [])
        slide_data = arr[slide_num - 1]
    
        if len(text_box_ids) != 2:
            print(f"⚠️ Slide {slide_num} mapping should have 2 box IDs but found {len(text_box_ids)}")
            continue

        title_box_id, content_box_id = text_box_ids

        try:
            # Update title
            title_text = slide_data[0]
            update_text_of_textbox(presentation, slide_num, title_box_id, title_text)

            # Update content
            content_text = slide_data[1]
            if isinstance(content_text, list):
                content_text = "\n".join(str(line) for line in content_text)
            update_text_of_textbox(presentation, slide_num, content_box_id, content_text)
    
        except IndexError as e:
            print(f"⚠️ Slide {slide_num} content missing: {e}")


    # Save the updated presentation
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    pptx_filename = f"generated_presentation_{timestamp}.pptx"
    pptx_path = os.path.join("powerpoints", pptx_filename)

    presentation.save(pptx_path)
    print(f"✅ Presentation saved at: {pptx_path}")

def slide_generation_2(arr, template_type):
    # Load the base presentation
    presentation = pptx.Presentation(os.path.join("powerpoints", f"{template_type}.pptx"))

    def update_text_of_textbox_by_name(presentation, slide_num, shape_name, new_text):
        slide = presentation.slides[slide_num - 1]
        for shape in slide.shapes:
            if shape.name == shape_name and shape.has_text_frame:
                text_frame = shape.text_frame
                first_paragraph = text_frame.paragraphs[0]
                first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()

                # Preserve formatting
                font = first_run.font
                font_name, font_size = font.name, font.size
                font_bold, font_italic = font.bold, font.italic
                font_underline = font.underline
                try:
                    font_color = font.color.rgb
                except:
                    font_color = None

                # Clear and add new text
                text_frame.clear()
                run = text_frame.paragraphs[0].add_run()
                run.text = str(new_text)

                # Reapply formatting
                run.font.name = font_name
                run.font.size = font_size
                run.font.bold = font_bold
                run.font.italic = font_italic
                run.font.underline = font_underline
                if font_color:
                    try:
                        run.font.color.rgb = font_color
                    except Exception as e:
                        print(f"⚠️ Font color error on Slide {slide_num}, Shape '{shape_name}': {e}")
                print(f"📝 Updated Slide {slide_num}, Shape '{shape_name}' → '{new_text[:30]}...'")
                return
        print(f"❌ Could not find shape '{shape_name}' on Slide {slide_num}")

    # Template-specific mappings
    template_mappings = {
        "template2": {
            1: ["TextBox 12", "TextBox 17"],
            2: ["TextBox 7", "TextBox 2"],
            3: ["TextBox 2", "TextBox 4"],
            4: ["TextBox 8", "TextBox 15"],
            5: ["TextBox 2", "TextBox 3"],
            6: ["TextBox 2", "TextBox 16"],
            7: ["TextBox 2", "TextBox 4"],
            8: ["TextBox 2", "TextBox 5"],
            9: ["TextBox 12", "TextBox 14"]
        }
    }

    if template_type not in template_mappings:
        raise ValueError(f"❌ Unknown template type: {template_type}")

    mapping = template_mappings[template_type]

    for slide_num, (title_shape, content_shape) in mapping.items():
        try:
            title, content = arr[slide_num - 1]
            update_text_of_textbox_by_name(presentation, slide_num, title_shape, title)
            update_text_of_textbox_by_name(presentation, slide_num, content_shape, content)
        except IndexError as e:
            print(f"⚠️ Slide {slide_num} content missing in array: {e}")
        except Exception as e:
            print(f"❌ Unexpected error on Slide {slide_num}: {e}")

    # Save final presentation
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    pptx_filename = f"generated_presentation_{timestamp}.pptx"
    pptx_path = os.path.join("powerpoints", pptx_filename)
    presentation.save(pptx_path)
    print(f"✅ Presentation saved at: {pptx_path}")

